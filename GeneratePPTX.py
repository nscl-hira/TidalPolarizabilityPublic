from pptx import Presentation
from pptx.util import Inches

left = Inches(1)
top = Inches(1.55)
height = Inches(5.5)
width = Inches(14)

def ImageOnlySlide(prs, title, img):
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.width = Inches(14)
    title_shape.height = Inches(1.5)
    title_shape.top = Inches(0.25)
    title_shape.left = Inches(1)

    shapes.title.text = title
    slide.shapes.add_picture(img, left, top, height=Inches(8))

def AddDescriptionSlide(prs, title, content):
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.width = Inches(14)
    title_shape.height = Inches(1.5)
    title_shape.top = Inches(0.25)
    title_shape.left = Inches(1)
    body_shape = shapes.placeholders[1]
    body_shape.width = Inches(14)
    body_shape.height = Inches(8)
    body_shape.top = Inches(2)
    body_shape.left = Inches(1)

    title_shape.text = title
    tf = body_shape.text_frame
    for cont in content:
        p = tf.add_paragraph()
        p.text = cont
        p.level = 0

def CreateFirstSlide(title_text, subtitle_text):
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(11)
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text
    subtitle.text = subtitle_text
    return prs

